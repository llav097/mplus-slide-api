from fastapi import FastAPI, Form
from fastapi.responses import FileResponse
from pptx import Presentation
import os
import uuid

app = FastAPI()

@app.post("/generate-pptx")
async def generate_pptx(slide_text: str = Form(...)):
    prs = Presentation()
    current_slide = None

    for line in slide_text.splitlines():
        if line.startswith("Slide"):
            slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(slide_layout)
            title = current_slide.shapes.title
            body = current_slide.placeholders[1]
            title.text = line.split(":")[1].strip()
            body.text = ""
        elif line.startswith("-") and current_slide:
            body.text += line + "\n"

    filename = f"slides_{uuid.uuid4().hex}.pptx"
    filepath = f"./{filename}"
    prs.save(filepath)

    return FileResponse(filepath, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)